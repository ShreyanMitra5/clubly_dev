import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openai import OpenAI
from dotenv import load_dotenv
import json

# Load environment variables
load_dotenv()

# Configure OpenRouter OpenAI client
client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=os.getenv('OPENROUTER_API_KEY'),
)

MODEL_NAME = "deepseek/deepseek-r1-0528-qwen3-8b:free"

# Slide templates
SLIDE_TEMPLATES = {
    "title": {
        "layout": 0,  # Title slide layout
        "title_size": Pt(44),
        "subtitle_size": Pt(32)
    },
    "content": {
        "layout": 1,  # Content slide layout
        "title_size": Pt(36),
        "content_size": Pt(24)
    },
    "code": {
        "layout": 6,  # Blank slide layout
        "title_size": Pt(36),
        "code_size": Pt(18)
    }
}

def create_slide(prs, template_type, title, content=None, subtitle=None, speaker_notes=None):
    """Create a slide with the specified template type and content."""
    template = SLIDE_TEMPLATES[template_type]
    slide = prs.slides.add_slide(prs.slide_layouts[template["layout"]])
    
    # Add title
    if template_type == "title":
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = subtitle
        title_shape.text_frame.paragraphs[0].font.size = template["title_size"]
        subtitle_shape.text_frame.paragraphs[0].font.size = template["subtitle_size"]
    else:
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = template["title_size"]
    
    # Add content based on template type
    if template_type == "content" and content:
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        for point in content:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = template["content_size"]
    
    elif template_type == "code" and content:
        # Add code in a textbox
        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        code_frame = code_box.text_frame
        code_frame.text = content
        code_frame.paragraphs[0].font.name = 'Courier New'
        code_frame.paragraphs[0].font.size = template["code_size"]
    
    # Add speaker notes if provided
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_notes
    
    return slide

def parse_title_content(content):
    """Parse title content with fallback options."""
    try:
        # Try the expected format first
        if "TITLE:" in content and "SUBTITLE:" in content:
            title = content.split("TITLE:")[1].split("SUBTITLE:")[0].strip()
            subtitle = content.split("SUBTITLE:")[1].strip()
        else:
            # Fallback: try to extract title and subtitle from the content
            lines = content.strip().split('\n')
            title = lines[0].strip()
            subtitle = lines[1].strip() if len(lines) > 1 else f"Week {week_number} - {topic}"
        
        return title, subtitle
    except Exception as e:
        print(f"Error parsing title content: {e}")
        print("Raw content:", content)
        # Return default values if parsing fails
        return f"{club_type} - Week {week_number}", f"Topic: {topic}"

def parse_slide_content(content):
    """Parse slide content with fallback options."""
    try:
        # Try the expected format first
        if "TITLE:" in content and "BULLETS:" in content and "NOTES:" in content:
            title = content.split("TITLE:")[1].split("BULLETS:")[0].strip()
            bullets_text = content.split("BULLETS:")[1].split("NOTES:")[0]
            notes = content.split("NOTES:")[1].strip()
        else:
            # Fallback: try to extract content from the text
            lines = content.strip().split('\n')
            title = lines[0].strip()
            bullets = []
            notes = ""
            in_bullets = False
            in_notes = False
            
            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue
                if line.startswith('-') or line.startswith('•'):
                    bullets.append(line.lstrip('-• '))
                    in_bullets = True
                elif in_bullets and not in_notes:
                    notes = line
                    in_notes = True
                elif in_notes:
                    notes += "\n" + line
            
            if not bullets:
                bullets = ["Key point 1", "Key point 2", "Key point 3"]
            if not notes:
                notes = "Speaker notes for this slide"
        
        # Extract bullets if not already done
        if "bullets" not in locals():
            bullets = [b.strip() for b in bullets_text.split("-") if b.strip()]
        
        return {
            "title": title,
            "content": bullets,
            "speaker_notes": notes
        }
    except Exception as e:
        print(f"Error parsing slide content: {e}")
        print("Raw content:", content)
        # Return default values if parsing fails
        return {
            "title": "Slide Title",
            "content": ["Point 1", "Point 2", "Point 3"],
            "speaker_notes": "Speaker notes for this slide"
        }

def generate_slide_content(club_type, topic, week_number):
    """Generate content for each slide using the LLM."""
    prompts = [
        {
            "role": "system",
            "content": "You are a helpful assistant that generates educational content for club presentations. Provide responses in a clear, structured format."
        },
        {
            "role": "user",
            "content": f"""Generate a title and subtitle for a {club_type} presentation about {topic} for week {week_number}.
            Format: TITLE: [title]
            SUBTITLE: [subtitle]"""
        }
    ]
    
    # Get title slide content
    completion = client.chat.completions.create(
        model=MODEL_NAME,
        messages=prompts,
        temperature=0.7,
        max_tokens=100
    )
    title_content = completion.choices[0].message.content
    
    # Parse title content
    title, subtitle = parse_title_content(title_content)
    
    # Generate content slides
    slides = []
    topics = [
        f"Introduction to {topic}",
        f"Key Concepts of {topic}",
        f"Examples and Applications of {topic}",
        f"Hands-on Practice with {topic}",
        f"Summary and Next Steps"
    ]
    
    for slide_topic in topics:
        prompts = [
            {
                "role": "system",
                "content": "You are a helpful assistant that generates educational content for club presentations. Provide responses in a clear, structured format."
            },
            {
                "role": "user",
                "content": f"""Generate content for a slide about {slide_topic} for a {club_type} presentation.
                Include:
                1. A title
                2. 3-5 bullet points
                3. Speaker notes
                
                Format:
                TITLE: [title]
                BULLETS:
                - [bullet 1]
                - [bullet 2]
                - [bullet 3]
                NOTES: [speaker notes]"""
            }
        ]
        
        completion = client.chat.completions.create(
            model=MODEL_NAME,
            messages=prompts,
            temperature=0.7,
            max_tokens=300
        )
        
        content = completion.choices[0].message.content
        slide_data = parse_slide_content(content)
        slides.append(slide_data)
    
    return {
        "title_slide": {
            "title": title,
            "subtitle": subtitle
        },
        "content_slides": slides
    }

def generate_slides(club_type: str, topic: str, week_number: int):
    """Generate a PowerPoint presentation for a club session."""
    # Create presentation
    prs = Presentation()
    
    # Generate content
    content = generate_slide_content(club_type, topic, week_number)
    
    # Create title slide
    create_slide(
        prs,
        "title",
        content["title_slide"]["title"],
        subtitle=content["title_slide"]["subtitle"]
    )
    
    # Create content slides
    for slide_data in content["content_slides"]:
        create_slide(
            prs,
            "content",
            slide_data["title"],
            content=slide_data["content"],
            speaker_notes=slide_data["speaker_notes"]
        )
    
    # Save the presentation
    filename = f"{club_type}_Week{week_number}.pptx"
    prs.save(filename)
    return filename

if __name__ == "__main__":
    # Example usage
    club_type = "Python Club"
    topic = "Introduction to Variables and Data Types"
    week_number = 1
    
    filename = generate_slides(club_type, topic, week_number)
    print(f"Generated presentation: {filename}") 