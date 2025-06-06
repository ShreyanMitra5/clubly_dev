import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import requests
from PIL import Image
from io import BytesIO
from openai import OpenAI
from dotenv import load_dotenv
import tempfile
from serpapi import GoogleSearch
import textwrap

# Load environment variables
load_dotenv()

# Slide themes and styles
THEMES = {
    "modern": {
        "background_color": RGBColor(250, 250, 250),  # Light gray background
        "title_color": RGBColor(33, 33, 33),         # Dark gray for titles
        "text_color": RGBColor(66, 66, 66),          # Medium gray for body text
        "accent_color": RGBColor(0, 122, 255),       # iOS blue for accents
        "font_family": "Montserrat",                 # Modern sans-serif font
        "title_font_size": Pt(40),
        "subtitle_font_size": Pt(28),
        "body_font_size": Pt(24),
        "bullet_font_size": Pt(22),
        "margin": Inches(1),
        "image_width": Inches(4),
        "image_height": Inches(3)
    },
    "dark": {
        "background_color": RGBColor(33, 33, 33),    # Dark background
        "title_color": RGBColor(255, 255, 255),      # White for titles
        "text_color": RGBColor(200, 200, 200),       # Light gray for body text
        "accent_color": RGBColor(0, 199, 190),       # Teal for accents
        "font_family": "Roboto",                     # Clean sans-serif font
        "title_font_size": Pt(40),
        "subtitle_font_size": Pt(28),
        "body_font_size": Pt(24),
        "bullet_font_size": Pt(22),
        "margin": Inches(1),
        "image_width": Inches(4),
        "image_height": Inches(3)
    },
    "nature": {
        "background_color": RGBColor(245, 245, 240), # Soft cream background
        "title_color": RGBColor(46, 64, 46),         # Forest green for titles
        "text_color": RGBColor(64, 64, 64),          # Dark gray for body text
        "accent_color": RGBColor(76, 175, 80),       # Green for accents
        "font_family": "Open Sans",                  # Friendly sans-serif font
        "title_font_size": Pt(40),
        "subtitle_font_size": Pt(28),
        "body_font_size": Pt(24),
        "bullet_font_size": Pt(22),
        "margin": Inches(1),
        "image_width": Inches(4),
        "image_height": Inches(3)
    },
    "coding": {
        "background_color": RGBColor(30, 32, 34),
        "title_color": RGBColor(80, 250, 123),
        "text_color": RGBColor(248, 248, 242),
        "accent_color": RGBColor(139, 233, 253),
        "font_family": "Consolas",  # Monospaced font for code feel
        "title_font_size": Pt(40),
        "subtitle_font_size": Pt(28),
        "body_font_size": Pt(22),
        "bullet_font_size": Pt(20),
        "margin": Inches(0.7),
        "image_width": Inches(3.5),
        "image_height": Inches(2.7)
    },
    "academic": {
        "background_color": RGBColor(255, 255, 255), # White background
        "title_color": RGBColor(20, 20, 140),       # Dark blue for titles
        "text_color": RGBColor(50, 50, 50),        # Dark gray for body text
        "accent_color": RGBColor(100, 100, 255),    # Lighter blue accent
        "font_family": "Times New Roman",           # Serif font for formal feel
        "title_font_size": Pt(44),
        "subtitle_font_size": Pt(30),
        "body_font_size": Pt(26),
        "bullet_font_size": Pt(24),
        "margin": Inches(1.2),
        "image_width": Inches(3.8),
        "image_height": Inches(2.8)
    },
    "creative": {
        "background_color": RGBColor(255, 240, 230), # Light peach background
        "title_color": RGBColor(150, 50, 0),        # Warm brown for titles
        "text_color": RGBColor(80, 40, 0),         # Darker brown for text
        "accent_color": RGBColor(255, 150, 50),     # Orange accent
        "font_family": "Georgia",                 # Serif font with character
        "title_font_size": Pt(42),
        "subtitle_font_size": Pt(28),
        "body_font_size": Pt(23),
        "bullet_font_size": Pt(21),
        "margin": Inches(0.9),
        "image_width": Inches(4.2),
        "image_height": Inches(3.1)
    }
}

def get_image_url(query, serpapi_key):
    """Get a relevant image URL using SerpAPI."""
    try:
        params = {
            "engine": "google",
            "q": query,
            "tbm": "isch",  # Image search
            "api_key": serpapi_key
        }
        
        search = GoogleSearch(params)
        results = search.get_dict()
        
        # Get the first image URL from the results
        if "images_results" in results and len(results["images_results"]) > 0:
            return results["images_results"][0]["original"]
        return None
    except Exception as e:
        print(f"Error getting image from SerpAPI: {e}")
        return None

def download_and_resize_image(image_url, width=800, height=600):
    """Download and resize an image for use in slides."""
    if not image_url:
        print("Warning: No image URL provided for download.")
        return None
    try:
        response = requests.get(image_url, timeout=10) # Added a timeout
        response.raise_for_status() # Raise an HTTPError for bad responses (4xx or 5xx)
        
        img = Image.open(BytesIO(response.content))
        # Convert to RGB if not already (fixes 'cannot write mode P as JPEG')
        if img.mode != 'RGB':
            img = img.convert('RGB')
        img = img.resize((width, height), Image.Resampling.LANCZOS)
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
        img.save(temp_file.name)
        print(f"Successfully downloaded and resized image from {image_url}")
        return temp_file.name
    except requests.exceptions.RequestException as e:
        print(f"Error downloading image from {image_url}: {e}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during image processing for {image_url}: {e}")
        return None

def create_title_slide(prs, title, subtitle, image_path, theme="modern"):
    """Create a title slide with background image."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEMES[theme]["background_color"]
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = THEMES[theme]["title_font_size"]
    title_frame.paragraphs[0].font.color.rgb = THEMES[theme]["title_color"]
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = THEMES[theme]["font_family"]
    
    # Add subtitle
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.paragraphs[0].font.size = THEMES[theme]["subtitle_font_size"]
    subtitle_frame.paragraphs[0].font.color.rgb = THEMES[theme]["text_color"]
    subtitle_frame.paragraphs[0].font.name = THEMES[theme]["font_family"]
    
    # Add image if available
    if image_path:
        # Position image on the right side
        left = prs.slide_width - THEMES[theme]["image_width"] - THEMES[theme]["margin"]
        top = THEMES[theme]["margin"]
        pic = slide.shapes.add_picture(
            image_path,
            left,
            top,
            width=THEMES[theme]["image_width"],
            height=THEMES[theme]["image_height"]
        )

def fit_text_to_box(text, max_chars=350):
    """Shorten or summarize text to fit in the slide."""
    if len(text) > max_chars:
        return textwrap.shorten(text, width=max_chars, placeholder="...")
    return text

def create_content_slide(prs, title, content, image_path, theme="modern", notes=None):
    """Create a content slide with bullet points and optional image, ensuring no overlap and text fits."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEMES[theme]["background_color"]
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = THEMES[theme]["title_font_size"]
    title_frame.paragraphs[0].font.color.rgb = THEMES[theme]["title_color"]
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = THEMES[theme]["font_family"]
    # Calculate available space
    margin = THEMES[theme]["margin"]
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    image_width = THEMES[theme]["image_width"] if image_path else 0
    image_height = THEMES[theme]["image_height"] if image_path else 0
    # Adjust text box width based on image presence
    # Increase text box width slightly when image is present for better fitting
    text_width = slide_width - image_width - (margin * 2.5) # Reduced right margin if image present
    # Ensure a minimum text box width even without an image
    if image_width == 0:
        text_width = slide_width - (2 * margin)
    content_shape = slide.placeholders[1]
    # Position text box on the left side, below title area
    text_box_left = int(margin)
    text_box_top = int(margin * 2.5) # Ensure it starts below title + a little space
    # Calculate text box height, ensuring it doesn't go past the bottom margin
    text_box_height = int(slide_height - text_box_top - margin)
    
    content_shape.left = text_box_left
    content_shape.top = text_box_top
    content_shape.width = int(text_width)
    content_shape.height = text_box_height
    
    tf = content_shape.text_frame
    tf.clear()
    tf.word_wrap = True
    
    # Process bullet points: split, clean, and filter out empty results
    # This is the critical part for handling unexpected newlines/whitespace
    raw_lines = content.split('\n')
    processed_bullet_lines = []
    for line in raw_lines:
        cleaned_line = line.strip().lstrip('-•* ')
        if cleaned_line: # Only add non-empty lines after cleaning
            processed_bullet_lines.append(cleaned_line)
            
    # Try to fit text with decreasing font size if needed
    min_font_size = 14
    default_font_size = THEMES[theme]["bullet_font_size"].pt
    font_size = default_font_size
    fits = False
    
    # Max attempts for font size reduction loop
    max_font_attempts = int(default_font_size - min_font_size) + 1
    for attempt in range(max_font_attempts):
        current_font_size = default_font_size - attempt
        if current_font_size < min_font_size:
            current_font_size = min_font_size
            
        tf.clear()
        # Add bullets with current font size
        for bullet in processed_bullet_lines:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(current_font_size)
            p.font.color.rgb = THEMES[theme]["text_color"]
            p.font.name = THEMES[theme]["font_family"]
            p.level = 0
            
        # Check if text fits within the text box height using a heuristic
        # The bounding box of the text frame after adding text is the best available check
        # If the text frame's height is less than or equal to the shape's height, it should fit
        # Adding a small buffer (e.g., 100000 EMU) for safety
        buffer_emu = 100000 # Approximately 0.1 inches
        if tf._extents[1] <= content_shape.height - buffer_emu:
             fits = True
             break
            
        if current_font_size == min_font_size: # Stop if we reached min font size and still doesn't fit
             break
            
    # If text still doesn't fit after reducing font size, log a warning.
    # At this point, the content might be fundamentally too long or the rendering is tricky.
    if not fits:
         print(f"Warning: Text overflowed on slide '{title}' even after reducing font size and cleaning newlines. Consider simplifying content.")
         # No further AI rephrasing during presentation generation due to complexity and potential infinite loops.
    
    # Add image if available
    if image_path:
        # Position image on the right side, below the title area
        image_left = slide_width - image_width - margin
        image_top = margin * 2.5 # Align top roughly with text box top
        # Ensure image does not go below the slide bottom margin
        if image_top + image_height > slide_height - margin:
             image_top = slide_height - image_height - margin
        # Ensure image does not go above the text box top if text box is very high
        if image_top < text_box_top:
            image_top = text_box_top
            
        slide.shapes.add_picture(
            image_path,
            image_left,
            image_top,
            width=image_width,
            height=image_height # Use theme defined image height
        )
    
    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

def should_add_image(title, content):
    """Determine if an image would be valuable for this slide."""
    # Keywords that suggest an image would be helpful
    image_keywords = [
        'example', 'illustration', 'diagram', 'visual', 'picture',
        'show', 'demonstrate', 'display', 'view', 'see'
    ]
    
    # Check title and content for image-related keywords
    text = (title + ' ' + content).lower()
    return any(keyword in text for keyword in image_keywords)

def generate_slide_content(club_type, topic, week_number, openrouter_key):
    """Generate content for each slide using the LLM."""
    # Generate title slide content (no image)
    prompts = [
        {
            "role": "system",
            "content": f"You are a helpful assistant that generates educational content for {club_type} presentations. Provide responses in a clear, structured format."
        },
        {
            "role": "user",
            "content": f"""Generate a title and subtitle for a {club_type} presentation about {topic} for week {week_number}.
            Format: TITLE: [title]\nSUBTITLE: [subtitle]"""
        }
    ]
    # Configure OpenRouter OpenAI client for this function call
    client = OpenAI(
        base_url="https://openrouter.ai/api/v1",
        api_key=openrouter_key,
    )
    
    completion = client.chat.completions.create(
        messages=prompts,
        model="google/gemini-2.5-flash-preview-05-20",
        temperature=0.7,
        max_tokens=100,
        extra_headers={
            "HTTP-Referer": "https://openrouter.ai/",
            "X-Title": "ClubAI-PPT-Generator"
        },
        extra_body={}
    )
    title_content = completion.choices[0].message.content
    print("API Response (Title Slide):", title_content)  # Debug print
    
    # Robust parsing for title and subtitle
    title_start = title_content.find("TITLE:")
    subtitle_start = title_content.find("SUBTITLE:")

    if title_start != -1:
        title_end = subtitle_start if subtitle_start != -1 else len(title_content)
        title = title_content[title_start + len("TITLE:"):title_end].strip()
    
    if subtitle_start != -1:
        subtitle = title_content[subtitle_start + len("SUBTITLE:"):].strip()
        
    # Generate content slides
    slides = []
    topics = [
        f"Introduction to {topic}",
        f"Key Concepts of {topic}",
        f"Examples and Applications of {topic}",
        f"Hands-on Practice with {topic}",
        f"Summary and Next Steps"
    ]
    print(f"--- Starting content slide generation loop for {len(topics)} topics: {topics}") # Debug print before loop
    used_image_terms = set()
    for slide_topic in topics:
        print(f"--- Processing topic: {slide_topic}") # Debug print at start of loop
        slide_data = None
        for attempt in range(3):  # Try up to 3 times to get text that fits
            try:
                prompts = [
                    {
                        "role": "system",
                        "content": "You are a helpful assistant that generates educational content for club presentations. Provide responses in a clear, structured format."
                    },
                    {
                        "role": "user",
                        "content": f"""Generate content for a slide about {slide_topic} for a {club_type} presentation.
                        Include:\n1. A title\n2. 3-5 concise bullet points (each bullet should be short and fit easily on one line, with NO extra blank lines between bullets)\n3. Speaker notes (Write these in a friendly, conversational, and encouraging tone, as if you are presenting to a high school club. Include practical tips for explaining the slide content and engaging the audience, such as suggesting demos, asking questions, or highlighting key takeaways. Focus on helping the presenter deliver the information effectively.)\n4. A relevant, unique, and universal programming-related image search term (not just Python, avoid repeating previous terms)\n                        Format:\nTITLE: [title]\nBULLETS:\n- [bullet 1]\n- [bullet 2]\n- [bullet 3]\nNOTES: [speaker notes]\nIMAGE: [image search term]"""
                    }
                ]
                # Configure OpenRouter OpenAI client for this attempt
                client = OpenAI(
                    base_url="https://openrouter.ai/api/v1",
                    api_key=openrouter_key,
                )

                completion = client.chat.completions.create(
                    model="google/gemini-2.5-flash-preview-05-20",
                    messages=prompts,
                    temperature=0.7,
                    max_tokens=300,
                    extra_headers={
                        "HTTP-Referer": "https://openrouter.ai/",
                        "X-Title": "ClubAI-PPT-Generator"
                    },
                    extra_body={}
                )
                content = completion.choices[0].message.content
                print(f"API Response for {slide_topic}:", content)  # Debug print
                slide_data = parse_slide_content(content)
                # Check if text fits (roughly, by line count and char count)
                if (len(slide_data["content"].split('\n')) <= 6 and
                    all(len(line) <= 90 for line in slide_data["content"].split('\n'))):
                    # Check for unique image term
                    if slide_data["image_term"] in used_image_terms:
                        continue  # Ask AI again for a new image term
                    used_image_terms.add(slide_data["image_term"])
                    break  # Accept this version
            except Exception as e:
                print(f"Error generating content for {slide_topic}: {e} - Raw content was:\n---\n{content if 'content' in locals() else 'N/A'}\n---") # Added debug print
            if not slide_data:
                # Fallback content if API call fails or parsing/fitting fails after retries
                slide_data = {
                    "title": slide_topic,
                    "content": f"Key points about {slide_topic}.", # Simplified fallback
                    "notes": f"Speaker notes for {slide_topic}.",
                    "image_term": f"{club_type} {slide_topic} illustration"
                }
        # Append slide data after attempts
        slides.append(slide_data)
        print(f"--- Appended slide data for topic: {slide_topic}") # Debug print after append
    
    print(f"Generated {len(slides)} content slides.") # Debug print
    print(f"Content slides data: {slides}") # Debug print
    
    return {
        "title_slide": {
            "title": title,
            "subtitle": subtitle
        },
        "content_slides": slides
    }

def parse_slide_content(content):
    """Parse the content generated by the LLM more robustly."""
    print(f"Attempting to parse content:\n---\n{content}\n---") # Debug print raw content
    slide_data = {
        "title": "Slide Title (Parsing Failed)", # Fallback title
        "content": "- Failed to parse content\n- Please check API response", # Fallback content
        "notes": "Parsing failed for this slide's content.", # Fallback notes
        "image_term": "error" # Fallback image term
    }
    
    try:
        # Use find and slicing to extract sections, making it more robust
        title_start = content.find("TITLE:")
        bullets_start = content.find("BULLETS:")
        notes_start = content.find("NOTES:")
        image_start = content.find("IMAGE:")
        
        if title_start != -1:
            title_end = bullets_start if bullets_start != -1 else len(content)
            slide_data["title"] = content[title_start + len("TITLE:"):title_end].strip()
            
        if bullets_start != -1:
            bullets_end = notes_start if notes_start != -1 else len(content)
            bullets_text = content[bullets_start + len("BULLETS:"):bullets_end].strip()
            
            # Clean up bullet points: split by lines, filter empty, remove leading markers
            bullets = []
            for line in bullets_text.split('\n'):
                cleaned_line = line.strip().lstrip('-•* ')
                if cleaned_line:
                    bullets.append(cleaned_line)
            slide_data["content"] = '\n'.join(bullets)
            
        if notes_start != -1:
            notes_end = image_start if image_start != -1 else len(content)
            slide_data["notes"] = content[notes_start + len("NOTES:"):notes_end].strip()
            
        if image_start != -1:
            slide_data["image_term"] = content[image_start + len("IMAGE:"):].strip()
        else:
            # Fallback image term if not provided by AI
            slide_data["image_term"] = slide_data.get("title", "chemistry illustration") 
            
    except Exception as e:
        print(f"Error parsing slide content: {e} - Raw content was:\n---\n{content}\n---")
        # Fallback values are already set at the beginning of the function
        
    return slide_data

def generate_presentation(club_type: str, topic: str, week_number: int, theme: str = "modern", openrouter_key: str = None, serpapi_key: str = None):
    """Generate a PowerPoint presentation with AI content and images."""
    prs = Presentation()
    content = generate_slide_content(club_type, topic, week_number, openrouter_key)
    # Create title slide (no image)
    create_title_slide(
        prs,
        content['title_slide']['title'],
        content['title_slide']['subtitle'],
        None,
        theme
    )
    # Create content slides
    for slide_data in content['content_slides']:
        image_url = get_image_url(slide_data['image_term'], serpapi_key)
        image_path = None
        if image_url:
            image_path = download_and_resize_image(image_url)
        create_content_slide(
            prs,
            slide_data['title'],
            slide_data['content'],
            image_path,
            theme,
            notes=slide_data.get('notes')
        )
        if image_path:
            os.unlink(image_path)
    filename = f"{club_type}_Week{week_number}.pptx"
    prs.save(filename)
    return filename

if __name__ == "__main__":
    # Get input from the user
    load_dotenv() # Load keys from .env when running standalone
    OPENROUTER_API_KEY = os.getenv('OPENROUTER_API_KEY')
    SERPAPI_KEY = os.getenv('SERPAPI_KEY')

    if not OPENROUTER_API_KEY or not SERPAPI_KEY:
        print("Please ensure OPENROUTER_API_KEY and SERPAPI_KEY are set in your .env file.")
        exit()

    # Now call generate_presentation with keys loaded from .env for standalone testing
    # (In the FastAPI app, these keys will come from the request)
    club_type = input("Enter the type of club (e.g., Biology Club, History Club): ")
    topic = input("Enter the topic of the presentation (e.g., Photosynthesis, World War II): ")
    week_number_str = input("Enter the week number: ")
    try:
        week_number = int(week_number_str)
    except ValueError:
        print("Invalid week number. Please enter a number.")
        exit()
    
    # List available themes
    print("\nAvailable themes:")
    for theme_name in THEMES.keys():
        print(f"- {theme_name}")
    
    theme = input(f"\nSelect a theme ({list(THEMES.keys())[0]}): ")
    if theme not in THEMES:
        print(f"Invalid theme selected. Using default theme: {list(THEMES.keys())[0]}")
        theme = list(THEMES.keys())[0] # Default to the first theme

    # Generate the presentation
    filename = generate_presentation(club_type, topic, week_number, theme, OPENROUTER_API_KEY, SERPAPI_KEY)
    print(f"Generated presentation: {filename}") 